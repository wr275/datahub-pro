import os
import io
import json
import math
import httpx
import logging
from collections import Counter
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse, Response
from pydantic import BaseModel
from typing import Optional, Any

from database import get_db, DataFile, User
from sqlalchemy.orm import Session
from auth_utils import get_current_user
from routers.admin import log_usage
from routers.analytics import load_file_data

logger = logging.getLogger(__name__)

# Claude Haiku 4.5 list pricing (USD per million tokens):
#   input  $1.00 → 100 cents / 1,000,000 tokens
#   output $5.00 → 500 cents / 1,000,000 tokens
# Using integer math on cents × tokens keeps everything precise.
HAIKU_INPUT_CENTS_PER_MTOK = 100
HAIKU_OUTPUT_CENTS_PER_MTOK = 500


def _estimate_cost_cents(input_tokens: int, output_tokens: int) -> int:
    total_microcents = (
        input_tokens * HAIKU_INPUT_CENTS_PER_MTOK
        + output_tokens * HAIKU_OUTPUT_CENTS_PER_MTOK
    )
    # Round up so we never under-bill (1,000,000 tokens ≈ 1 cent floor).
    return (total_microcents + 999_999) // 1_000_000

router = APIRouter(prefix="/ai", tags=["ai"])


def _get_owned_file(file_id: str, user: User, db: Session) -> DataFile:
    """Return a DataFile that belongs to the caller's org, or raise 404.
    404 (not 403) so we don't leak existence of other orgs' file IDs."""
    f = db.query(DataFile).filter(
        DataFile.id == file_id,
        DataFile.organisation_id == user.organisation_id,
    ).first()
    if not f:
        raise HTTPException(status_code=404, detail="File not found")
    return f


def require_ai_enabled(current_user: User = Depends(get_current_user)) -> User:
    """Gate that refuses the call when the caller's org hasn't enabled the
    AI add-on. Returns a 403 with a machine-readable ``code`` so the frontend
    can render the upgrade CTA rather than a generic error toast.
    """
    org = current_user.organisation
    if not org or not bool(getattr(org, "ai_enabled", False)):
        raise HTTPException(
            status_code=403,
            detail={
                "code": "ai_disabled",
                "message": "AI features are not enabled for this workspace. An owner or admin can enable them in Organisation settings.",
            },
        )
    return current_user

ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
ANTHROPIC_URL = "https://api.anthropic.com/v1/messages"
MODEL = "claude-haiku-4-5-20251001"

SYSTEM_PROMPT = """You are a data analyst assistant embedded in a data hub application.
The user will ask questions about their uploaded data files (CSVs, spreadsheets, etc.).

When asked to show data as a table, output ONLY this raw JSON with no markdown, no code blocks, no backticks:
{"type":"table","summary":"brief description","columns":["col1","col2"],"rows":[{"col1":"val1","col2":"val2"}]}

When asked to show a chart or graph, output ONLY this raw JSON with no markdown, no code blocks, no backticks:
{"type":"chart","summary":"brief description","chart_type":"bar","chart_data":[{"name":"A","value":10}],"x_key":"name","y_keys":["value"],"title":"Chart Title"}

chart_type must be one of: bar, line, pie
CRITICAL: For table and chart responses, output ONLY the raw JSON object starting with { and ending with }. Never use markdown code blocks or backticks. Never add any text before or after the JSON.

For all other questions, respond with clear, concise plain text.
Be helpful, accurate, and data-focused. Keep responses brief unless detail is requested."""


class StreamRequest(BaseModel):
    file_id: str
    question: str


class GreetRequest(BaseModel):
    file_id: str
    filename: Optional[str] = None
    rows: Optional[int] = None
    cols: Optional[int] = None


class PromptRequest(BaseModel):
    file_id: str
    question: str


def build_context(file: DataFile, question: str) -> str:
    parts = [f"File: {file.filename}"]
    if file.row_count:
        parts.append(f"Rows: {file.row_count}")
    if file.column_count:
        parts.append(f"Columns: {file.column_count}")
    if file.columns_json:
        try:
            cols = json.loads(file.columns_json) if isinstance(file.columns_json, str) else file.columns_json
            parts.append(f"Column names: {', '.join(str(c) for c in cols[:20])}")
        except Exception:
            pass
    if file.file_content:
        try:
            raw = file.file_content
            preview = raw.decode("utf-8", errors="replace")[:3000] if isinstance(raw, bytes) else str(raw)[:3000]
            parts.append(f"\nData preview:\n{preview}")
        except Exception:
            pass
    parts.append(f"\nUser question: {question}")
    return "\n".join(parts)


@router.post("/stream")
async def stream_ai(
    req: StreamRequest,
    current_user: User = Depends(require_ai_enabled),
    db: Session = Depends(get_db),
):
    if not ANTHROPIC_API_KEY:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured")

    file = _get_owned_file(req.file_id, current_user, db)
    user_message = build_context(file, req.question)

    # Closure-local counters. Anthropic's SSE sends input_tokens in
    # `message_start` and the final output_tokens in `message_delta`
    # right before `message_stop`. We accumulate then log once the
    # stream completes (success or error).
    usage_state = {"input": 0, "output": 0}
    org_id = current_user.organisation_id
    user_id = current_user.id

    async def event_generator():
        try:
            async with httpx.AsyncClient(timeout=120) as client:
                async with client.stream(
                    "POST",
                    ANTHROPIC_URL,
                    headers={
                        "x-api-key": ANTHROPIC_API_KEY,
                        "anthropic-version": "2023-06-01",
                        "content-type": "application/json",
                    },
                    json={
                        "model": MODEL,
                        "max_tokens": 2048,
                        "stream": True,
                        "system": SYSTEM_PROMPT,
                        "messages": [{"role": "user", "content": user_message}],
                    },
                ) as response:
                    async for line in response.aiter_lines():
                        if not line.startswith("data: "):
                            continue
                        data = line[6:]
                        try:
                            parsed = json.loads(data)
                            event_type = parsed.get("type", "")
                            if event_type == "message_start":
                                u = parsed.get("message", {}).get("usage", {}) or {}
                                usage_state["input"] = int(u.get("input_tokens", 0) or 0)
                            elif event_type == "message_delta":
                                u = parsed.get("usage", {}) or {}
                                if u.get("output_tokens") is not None:
                                    usage_state["output"] = int(u.get("output_tokens", 0) or 0)
                            elif event_type == "content_block_delta":
                                delta = parsed.get("delta", {})
                                if delta.get("type") == "text_delta":
                                    text = delta.get("text", "")
                                    if text:
                                        yield f"data: {text}\n\n"
                            elif event_type == "message_stop":
                                yield "data: [DONE]\n\n"
                                break
                            elif event_type == "error":
                                error_msg = parsed.get("error", {}).get("message", "Unknown error")
                                yield f"data: Error: {error_msg}\n\n"
                                yield "data: [DONE]\n\n"
                                break
                        except Exception:
                            pass
        except Exception as e:
            yield f"data: Error: {str(e)}\n\n"
            yield "data: [DONE]\n\n"
        finally:
            # Best-effort usage log; errors here must never break the
            # SSE response that already went to the client.
            try:
                in_tok = usage_state["input"]
                out_tok = usage_state["output"]
                total = in_tok + out_tok
                if total > 0 and org_id:
                    log_usage(
                        db,
                        organisation_id=org_id,
                        user_id=user_id,
                        kind="ai_tokens",
                        quantity=total,
                        cost_cents=_estimate_cost_cents(in_tok, out_tok),
                        meta={
                            "model": MODEL,
                            "endpoint": "/ai/stream",
                            "input_tokens": in_tok,
                            "output_tokens": out_tok,
                            "file_id": req.file_id,
                        },
                    )
            except Exception as exc:
                logger.warning("usage logging failed in /ai/stream: %s", exc)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Access-Control-Allow-Origin": "*",
        },
    )


@router.post("/greet")
async def greet_file(
    req: GreetRequest,
    current_user: User = Depends(require_ai_enabled),
    db: Session = Depends(get_db),
):
    file = _get_owned_file(req.file_id, current_user, db)

    name = req.filename or file.filename or "your file"
    rows = req.rows or file.row_count or 0
    cols = req.cols or file.column_count or 0

    greeting = f"I've loaded **{name}**"
    if rows and cols:
        greeting += f" — {rows:,} rows, {cols} columns."
    elif rows:
        greeting += f" — {rows:,} rows."
    else:
        greeting += "."
    greeting += " Ask me anything about it!"

    return {"greeting": greeting}


@router.post("/prompt")
async def prompt_ai(
    req: PromptRequest,
    current_user: User = Depends(require_ai_enabled),
    db: Session = Depends(get_db),
):
    """Non-streaming fallback endpoint."""
    if not ANTHROPIC_API_KEY:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured")

    file = _get_owned_file(req.file_id, current_user, db)
    user_message = build_context(file, req.question)

    async with httpx.AsyncClient(timeout=60) as client:
        response = await client.post(
            ANTHROPIC_URL,
            headers={
                "x-api-key": ANTHROPIC_API_KEY,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": MODEL,
                "max_tokens": 2048,
                "system": SYSTEM_PROMPT,
                "messages": [{"role": "user", "content": user_message}],
            },
        )
        result = response.json()
        content = result.get("content", [{}])
        text = content[0].get("text", "") if content else ""

    # Usage logging: the non-streaming response includes `usage` at the top
    # level. Mirror the /stream accounting so both endpoints feed the same
    # ledger. Failures here must never surface to the caller.
    try:
        usage = result.get("usage", {}) or {}
        in_tok = int(usage.get("input_tokens", 0) or 0)
        out_tok = int(usage.get("output_tokens", 0) or 0)
        total = in_tok + out_tok
        if total > 0 and current_user.organisation_id:
            log_usage(
                db,
                organisation_id=current_user.organisation_id,
                user_id=current_user.id,
                kind="ai_tokens",
                quantity=total,
                cost_cents=_estimate_cost_cents(in_tok, out_tok),
                meta={
                    "model": MODEL,
                    "endpoint": "/ai/prompt",
                    "input_tokens": in_tok,
                    "output_tokens": out_tok,
                    "file_id": req.file_id,
                },
            )
    except Exception as exc:
        logger.warning("usage logging failed in /ai/prompt: %s", exc)

    return {"response": text}


# ---------------------------------------------------------------------------
# AI Insights 2.0 — stat-pack + LLM reasoning
# ---------------------------------------------------------------------------

INSIGHTS_SYSTEM_PROMPT = """You are a senior data analyst at a market research agency.
You are given a statistical summary ("stat pack") of a client dataset and must
produce an executive-quality insight brief. The user is non-technical — lead with
business meaning, not data mechanics.

Return ONLY raw JSON (no markdown, no code fences). Schema:

{
  "headline": "one punchy sentence (<=20 words) that names the most important thing in this data",
  "key_findings": [
    {
      "severity": "info | warn | alert",
      "category": "growth | concentration | trend | outlier | segmentation | opportunity",
      "title": "short bold-able title (<=8 words)",
      "detail": "2-3 sentences explaining what the number means in business terms, naming specific columns / segments / numbers from the stat pack",
      "metric": "the specific number that anchors the finding, e.g. '£1.2M, 34% of revenue'"
    }
  ],
  "data_quality": [
    {
      "severity": "info | warn | alert",
      "field": "column name or 'dataset'",
      "issue": "what's wrong or unusual",
      "recommendation": "concrete next action"
    }
  ],
  "opportunities": [
    {
      "title": "short verb-led suggestion (<=8 words)",
      "detail": "1-2 sentences on the business opportunity",
      "suggested_next_step": "the specific analysis or action the user should run next"
    }
  ]
}

Rules:
- 3-5 key_findings, 2-4 data_quality items, 2-4 opportunities.
- Every finding MUST cite at least one specific number or column from the stat pack.
- Prefer percentages + absolute numbers together ("£340k, 42% of total").
- If something is flagged with low_n=true, soften the language ("preliminary signal, small sample").
- Never invent columns, values, or categories that aren't in the stat pack.
- If the data is too sparse for a finding in a category, omit that item rather than padding.
"""


def _is_number(v: Any) -> bool:
    try:
        if v is None or v == "":
            return False
        float(v)
        return True
    except (ValueError, TypeError):
        return False


def _coerce_float(v: Any) -> Optional[float]:
    try:
        if v is None or v == "":
            return None
        return float(v)
    except (ValueError, TypeError):
        return None


def _looks_like_date(v: Any) -> bool:
    if not isinstance(v, str):
        return False
    s = v.strip()
    if len(s) < 6 or len(s) > 30:
        return False
    # Crude: YYYY-MM, YYYY-MM-DD, ISO, MM/DD/YYYY, dd-mm-yyyy
    import re
    return bool(re.match(r"^\d{4}[-/]\d{1,2}([-/]\d{1,2})?", s) or re.match(r"^\d{1,2}[-/]\d{1,2}[-/]\d{2,4}", s))


def _percentile(sorted_vals, p):
    if not sorted_vals:
        return None
    k = (len(sorted_vals) - 1) * p
    f = math.floor(k)
    c = math.ceil(k)
    if f == c:
        return sorted_vals[int(k)]
    return sorted_vals[f] * (c - k) + sorted_vals[c] * (k - f)


def build_stat_pack(rows: list, filename: str) -> dict:
    """Compute a rich statistical summary the LLM can reason over.

    The point of this function is to do the *math* on the server (so the LLM
    never has to chew through thousands of rows) and hand the LLM a compact
    JSON blob with all the signal it needs: shape, per-column stats, top
    categorical values with share, concentration ratios, outlier counts,
    date-range, and simple pairwise correlations for numeric columns.
    """
    if not rows:
        return {"filename": filename, "rows": 0, "columns": [], "empty": True}

    headers = list(rows[0].keys())
    n = len(rows)
    cols_out = []

    # Typing pass
    col_type: dict = {}
    col_numeric_vals: dict = {}
    col_text_vals: dict = {}
    for h in headers:
        nums = []
        texts = []
        nulls = 0
        date_hits = 0
        for r in rows:
            v = r.get(h)
            if v is None or v == "":
                nulls += 1
                continue
            fv = _coerce_float(v)
            if fv is not None:
                nums.append(fv)
            else:
                s = str(v)
                texts.append(s)
                if _looks_like_date(s):
                    date_hits += 1
        if nums and len(nums) >= max(3, int(n * 0.6)):
            col_type[h] = "numeric"
            col_numeric_vals[h] = nums
        elif date_hits and date_hits >= max(3, int(n * 0.5)):
            col_type[h] = "date"
            col_text_vals[h] = texts
        else:
            col_type[h] = "text"
            col_text_vals[h] = texts

        # Per-column output
        col_block = {"column": h, "type": col_type[h], "null_pct": round(nulls / n * 100, 1)}

        if col_type[h] == "numeric":
            vs = sorted(nums)
            s = sum(vs)
            mean = s / len(vs)
            var = sum((x - mean) ** 2 for x in vs) / len(vs) if len(vs) > 1 else 0
            std = math.sqrt(var)
            p25 = _percentile(vs, 0.25)
            p50 = _percentile(vs, 0.50)
            p75 = _percentile(vs, 0.75)
            # Outlier count: values more than 3 sd from the mean.
            out_hi = sum(1 for x in vs if std > 0 and x > mean + 3 * std)
            out_lo = sum(1 for x in vs if std > 0 and x < mean - 3 * std)
            zeros = sum(1 for x in vs if x == 0)
            negatives = sum(1 for x in vs if x < 0)
            col_block.update({
                "count": len(vs),
                "sum": round(s, 2),
                "mean": round(mean, 2),
                "median": round(p50, 2) if p50 is not None else None,
                "p25": round(p25, 2) if p25 is not None else None,
                "p75": round(p75, 2) if p75 is not None else None,
                "std": round(std, 2),
                "min": round(vs[0], 2),
                "max": round(vs[-1], 2),
                "zeros_pct": round(zeros / len(vs) * 100, 1),
                "negatives_pct": round(negatives / len(vs) * 100, 1),
                "outliers_high": out_hi,
                "outliers_low": out_lo,
            })
        elif col_type[h] == "text":
            counter = Counter(texts)
            uniq = len(counter)
            top = counter.most_common(7)
            top_share = [
                {"value": v, "count": c, "pct": round(c / len(texts) * 100, 1)}
                for v, c in top
            ]
            # Concentration: share of top-1 and top-3 relative to non-null.
            top1_share = round(top[0][1] / len(texts) * 100, 1) if top else 0
            top3_share = round(sum(c for _, c in top[:3]) / len(texts) * 100, 1) if top else 0
            col_block.update({
                "count": len(texts),
                "unique": uniq,
                "top_values": top_share,
                "top1_share_pct": top1_share,
                "top3_share_pct": top3_share,
                "is_likely_id": uniq >= max(int(len(texts) * 0.9), 5) and uniq > 20,
            })
        else:  # date
            try:
                parsed = sorted(t[:10] for t in texts if t)
                col_block.update({
                    "count": len(parsed),
                    "min_date": parsed[0] if parsed else None,
                    "max_date": parsed[-1] if parsed else None,
                    "unique_dates": len(set(parsed)),
                })
            except Exception:
                pass
        cols_out.append(col_block)

    # Pairwise Pearson correlations for numeric pairs (cap at 6 cols)
    correlations = []
    num_cols = [c for c in cols_out if c["type"] == "numeric"][:6]
    if len(num_cols) >= 2:
        # Re-extract aligned numeric arrays per row so indices line up.
        for i in range(len(num_cols)):
            for j in range(i + 1, len(num_cols)):
                a, b = num_cols[i]["column"], num_cols[j]["column"]
                xs, ys = [], []
                for r in rows:
                    fa = _coerce_float(r.get(a))
                    fb = _coerce_float(r.get(b))
                    if fa is not None and fb is not None:
                        xs.append(fa)
                        ys.append(fb)
                if len(xs) < 10:
                    continue
                mx = sum(xs) / len(xs)
                my = sum(ys) / len(ys)
                num = sum((x - mx) * (y - my) for x, y in zip(xs, ys))
                dx = math.sqrt(sum((x - mx) ** 2 for x in xs))
                dy = math.sqrt(sum((y - my) ** 2 for y in ys))
                if dx == 0 or dy == 0:
                    continue
                r = num / (dx * dy)
                if abs(r) >= 0.3:  # only report meaningful correlations
                    correlations.append({
                        "a": a,
                        "b": b,
                        "r": round(r, 2),
                        "n": len(xs),
                    })
        correlations.sort(key=lambda c: abs(c["r"]), reverse=True)
        correlations = correlations[:5]

    pack = {
        "filename": filename,
        "row_count": n,
        "column_count": len(headers),
        "columns": cols_out,
        "correlations": correlations,
        "low_n": n < 30,
    }
    return pack


class InsightsRequest(BaseModel):
    file_id: str


class FormulaRequest(BaseModel):
    file_id: str
    description: str  # plain-English "I want a formula that..."


class ChatMessage(BaseModel):
    role: str  # "user" or "assistant"
    content: str


class ChatRequest(BaseModel):
    file_id: str
    messages: list[ChatMessage]  # multi-turn conversation history


CHAT_SYSTEM_PROMPT = """You are a senior data analyst answering questions about
an uploaded dataset. The user will ask follow-up questions so remember previous
turns. When the user asks for data that is best presented as a table, respond
with ONLY this raw JSON (no markdown, no code fences):

{"type":"table","summary":"1-sentence framing","columns":["col1","col2"],"rows":[{"col1":"v","col2":"v"}]}

When the user asks for a chart or visualisation, respond with ONLY:

{"type":"chart","summary":"1-sentence framing","chart_type":"bar|line|pie","chart_data":[{"name":"A","value":10}],"x_key":"name","y_keys":["value"],"title":"Chart title"}

For everything else, respond with a concise, insight-led plain-text answer —
name specific numbers and columns from the data, and be confident but honest
about limits (e.g. low sample sizes). Avoid generic caveats."""


@router.post("/chat")
async def ai_chat(
    req: ChatRequest,
    current_user: User = Depends(require_ai_enabled),
    db: Session = Depends(get_db),
):
    """Multi-turn chat endpoint. The full message history is sent on every
    call so the model has conversation memory. We prepend a compact view of
    the file (columns + a small row sample) to the *first* user message so
    the model can ground its answers without ballooning token usage on
    every turn."""
    if not ANTHROPIC_API_KEY:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured")
    file = _get_owned_file(req.file_id, current_user, db)

    if not req.messages:
        raise HTTPException(status_code=400, detail="messages is empty")

    # Build a compact "dataset context" string that prefixes the first user turn.
    context_bits = [f"Dataset file: {file.filename}"]
    if file.row_count:
        context_bits.append(f"Rows: {file.row_count}")
    if file.column_count:
        context_bits.append(f"Columns: {file.column_count}")
    if file.columns_json:
        try:
            cols = json.loads(file.columns_json) if isinstance(file.columns_json, str) else file.columns_json
            context_bits.append(f"Column names: {', '.join(str(c) for c in cols[:25])}")
        except Exception:
            pass
    if file.file_content:
        try:
            raw = file.file_content
            preview = raw.decode("utf-8", errors="replace")[:2500] if isinstance(raw, bytes) else str(raw)[:2500]
            context_bits.append("\nData preview (first rows):\n" + preview)
        except Exception:
            pass
    context = "\n".join(context_bits)

    # Copy messages, prefix context onto the first user turn only.
    api_messages = []
    injected = False
    for m in req.messages:
        if m.role not in ("user", "assistant"):
            continue
        content = m.content
        if not injected and m.role == "user":
            content = context + "\n\nUser question: " + content
            injected = True
        api_messages.append({"role": m.role, "content": content})

    async with httpx.AsyncClient(timeout=90) as client:
        response = await client.post(
            ANTHROPIC_URL,
            headers={
                "x-api-key": ANTHROPIC_API_KEY,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": MODEL,
                "max_tokens": 2048,
                "system": CHAT_SYSTEM_PROMPT,
                "messages": api_messages,
            },
        )
        result = response.json()

    try:
        usage = result.get("usage", {}) or {}
        in_tok = int(usage.get("input_tokens", 0) or 0)
        out_tok = int(usage.get("output_tokens", 0) or 0)
        total = in_tok + out_tok
        if total > 0 and current_user.organisation_id:
            log_usage(
                db,
                organisation_id=current_user.organisation_id,
                user_id=current_user.id,
                kind="ai_tokens",
                quantity=total,
                cost_cents=_estimate_cost_cents(in_tok, out_tok),
                meta={"model": MODEL, "endpoint": "/ai/chat", "input_tokens": in_tok, "output_tokens": out_tok, "file_id": req.file_id, "turns": len(req.messages)},
            )
    except Exception as exc:
        logger.warning("usage logging failed in /ai/chat: %s", exc)

    content = result.get("content", [{}])
    text = content[0].get("text", "") if content else ""
    return {"response": text}


FORMULA_SYSTEM_PROMPT = """You generate mathjs-compatible formulas that operate
on a single row of a tabular dataset.

You will be given: (a) a list of available column names, (b) a plain-English
description from the user of the formula they want.

Return ONLY raw JSON (no markdown) of this shape:

{
  "expression": "mathjs expression using the exact column names as variables",
  "explanation": "1 sentence explaining what the formula does",
  "suggested_column_name": "Title_Case_With_Underscores",
  "confidence": "high | medium | low"
}

Rules for the expression:
- Use the column names as variables verbatim. If a column name contains spaces
  or special characters, it will be pre-mapped to a safe alias — reference
  columns as provided in the column list, not as strings.
- Supported functions: abs, round, floor, ceil, min, max, pow, sqrt, log,
  exp, mean, concat, substring, length, toLower, toUpper, trim, and the
  conditional: `cond ? a : b` or `if(cond, a, b)`.
- Supported operators: + - * / ^ ( ) > < >= <= == != and or not
- String literals use double quotes.
- Do NOT reference rows, ranges, or other rows — per-row only.
- If the request is ambiguous, pick the most likely interpretation and mark
  confidence "medium" or "low".
- If the request cannot be expressed per-row (e.g. needs a cross-row sum),
  set expression to "" and explanation should say why.
"""


@router.post("/formula")
async def ai_formula(
    req: FormulaRequest,
    current_user: User = Depends(require_ai_enabled),
    db: Session = Depends(get_db),
):
    """Plain-English → formula expression.

    The frontend sends the file_id + a description; we return a mathjs
    expression that references the dataset's columns by name. The frontend
    then evaluates it locally using mathjs against the actual rows.
    """
    if not ANTHROPIC_API_KEY:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured")

    file = _get_owned_file(req.file_id, current_user, db)
    # We only need column names (not full rows) — cheap path.
    cols = []
    if file.columns_json:
        try:
            cols = json.loads(file.columns_json) if isinstance(file.columns_json, str) else file.columns_json
        except Exception:
            cols = []
    if not cols:
        rows, _ = load_file_data(req.file_id, current_user.organisation_id, db)
        cols = list(rows[0].keys()) if rows else []

    user_message = (
        f"Available columns: {json.dumps(cols)}\n"
        f"User's request: {req.description}\n"
        "Return the formula JSON."
    )

    async with httpx.AsyncClient(timeout=45) as client:
        response = await client.post(
            ANTHROPIC_URL,
            headers={
                "x-api-key": ANTHROPIC_API_KEY,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": MODEL,
                "max_tokens": 800,
                "system": FORMULA_SYSTEM_PROMPT,
                "messages": [{"role": "user", "content": user_message}],
            },
        )
        result = response.json()

    try:
        usage = result.get("usage", {}) or {}
        in_tok = int(usage.get("input_tokens", 0) or 0)
        out_tok = int(usage.get("output_tokens", 0) or 0)
        total = in_tok + out_tok
        if total > 0 and current_user.organisation_id:
            log_usage(
                db,
                organisation_id=current_user.organisation_id,
                user_id=current_user.id,
                kind="ai_tokens",
                quantity=total,
                cost_cents=_estimate_cost_cents(in_tok, out_tok),
                meta={"model": MODEL, "endpoint": "/ai/formula", "input_tokens": in_tok, "output_tokens": out_tok, "file_id": req.file_id},
            )
    except Exception as exc:
        logger.warning("usage logging failed in /ai/formula: %s", exc)

    content = result.get("content", [{}])
    text = content[0].get("text", "") if content else ""
    raw = text.strip()
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1] if "\n" in raw else raw
        if raw.rstrip().endswith("```"):
            raw = raw.rstrip()[:-3]
    try:
        data = json.loads(raw)
    except Exception:
        data = {"expression": "", "explanation": "Model returned unparsable output.", "suggested_column_name": "Computed", "confidence": "low", "raw": text}
    return data


@router.post("/insights/{file_id}")
async def ai_insights(
    file_id: str,
    current_user: User = Depends(require_ai_enabled),
    db: Session = Depends(get_db),
):
    """Real LLM-backed insights.

    1. Load the file and compute a compact, information-dense stat pack.
    2. Send the stat pack (not the raw rows) to Claude Haiku with an
       executive-insight system prompt.
    3. Parse the JSON response and return it alongside the stat pack so the
       frontend can render cards plus supporting numbers.
    """
    if not ANTHROPIC_API_KEY:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured")

    _get_owned_file(file_id, current_user, db)  # authorisation check
    rows, filename = load_file_data(file_id, current_user.organisation_id, db)
    stat_pack = build_stat_pack(rows, filename)

    if stat_pack.get("empty"):
        return {
            "stat_pack": stat_pack,
            "insights": {
                "headline": "File is empty — no insights to generate.",
                "key_findings": [],
                "data_quality": [{
                    "severity": "alert",
                    "field": "dataset",
                    "issue": "The file contains no data rows.",
                    "recommendation": "Re-upload the file or check that the source export succeeded.",
                }],
                "opportunities": [],
            },
        }

    # Trim column payload so we stay within a sensible token budget for large files.
    trimmed = dict(stat_pack)
    trimmed["columns"] = stat_pack["columns"][:40]

    user_message = (
        "Stat pack for the uploaded dataset follows.\n\n"
        f"{json.dumps(trimmed, default=str)}\n\n"
        "Write the executive insight brief per the schema."
    )

    async with httpx.AsyncClient(timeout=90) as client:
        response = await client.post(
            ANTHROPIC_URL,
            headers={
                "x-api-key": ANTHROPIC_API_KEY,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": MODEL,
                "max_tokens": 2500,
                "system": INSIGHTS_SYSTEM_PROMPT,
                "messages": [{"role": "user", "content": user_message}],
            },
        )
        result = response.json()

    # Log token usage just like /ai/prompt does.
    try:
        usage = result.get("usage", {}) or {}
        in_tok = int(usage.get("input_tokens", 0) or 0)
        out_tok = int(usage.get("output_tokens", 0) or 0)
        total = in_tok + out_tok
        if total > 0 and current_user.organisation_id:
            log_usage(
                db,
                organisation_id=current_user.organisation_id,
                user_id=current_user.id,
                kind="ai_tokens",
                quantity=total,
                cost_cents=_estimate_cost_cents(in_tok, out_tok),
                meta={
                    "model": MODEL,
                    "endpoint": "/ai/insights",
                    "input_tokens": in_tok,
                    "output_tokens": out_tok,
                    "file_id": file_id,
                },
            )
    except Exception as exc:
        logger.warning("usage logging failed in /ai/insights: %s", exc)

    content = result.get("content", [{}])
    text = content[0].get("text", "") if content else ""

    # The model is instructed to return raw JSON, but guard against fenced output.
    raw = text.strip()
    if raw.startswith("```"):
        # Strip first fence line
        raw = raw.split("\n", 1)[1] if "\n" in raw else raw
        # Strip trailing fence
        if raw.rstrip().endswith("```"):
            raw = raw.rstrip()[:-3]
    try:
        insights = json.loads(raw)
    except Exception:
        # Fallback: hand the raw text back as a headline so the user isn't blocked.
        insights = {
            "headline": "AI returned an unexpected format. See raw text below.",
            "key_findings": [],
            "data_quality": [],
            "opportunities": [],
            "raw": text,
        }

    return {"stat_pack": stat_pack, "insights": insights}


# ---------------------------------------------------------------------------
# AI Narrative 2.0 — LLM-written prose tuned to audience / tone / length
# ---------------------------------------------------------------------------

NARRATIVE_SYSTEM_PROMPT = """You are a senior insight analyst writing prose
for an agency deck. You will receive:
- a statistical "stat pack" describing an uploaded dataset
- author preferences: audience, tone, length, optional focus column

Write a single continuous narrative — no lists, no headings, no markdown.
Ground every claim in specific numbers from the stat pack. Never invent data.

Tone guide:
- "executive": confident, outcome-led, decisive. No hedging.
- "analyst":   precise, includes sample sizes and caveats where relevant.
- "storyteller": uses narrative arc — opening hook, tension, resolution.
- "plain": simple words, no jargon, written for a non-technical reader.

Length guide:
- "short": 80-120 words, one paragraph.
- "medium": 180-240 words, two paragraphs.
- "long": 320-400 words, three paragraphs.

Audience framing:
- "client": address the reader as the client receiving the deliverable.
  Avoid internal-only language.
- "board": framed as briefing senior stakeholders. Emphasise "so what".
- "team": internal voice, team context, actionable.

If a focus_column is given, anchor the narrative around it; otherwise pick
the single most interesting numeric column from the stat pack (largest
concentration, biggest outlier, strongest correlation) and anchor there.

Return ONLY raw JSON:
{"narrative": "...", "headline": "short one-line headline under 12 words",
 "anchor_column": "name of the column you anchored on"}
"""


class NarrativeRequest(BaseModel):
    audience: str = "client"       # client | board | team
    tone: str = "executive"        # executive | analyst | storyteller | plain
    length: str = "medium"         # short | medium | long
    focus_column: Optional[str] = None


@router.post("/narrative/{file_id}")
async def ai_narrative(
    file_id: str,
    body: NarrativeRequest,
    current_user: User = Depends(require_ai_enabled),
    db: Session = Depends(get_db),
):
    """Generate a tuned written narrative about the dataset using the same
    stat pack the Report and Insights endpoints use. Keeps the page fast
    because the heavy statistics are computed server-side and only a compact
    JSON blob goes to the model.
    """
    if not ANTHROPIC_API_KEY:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured")
    _get_owned_file(file_id, current_user, db)
    rows, filename = load_file_data(file_id, current_user.organisation_id, db)
    stat_pack = build_stat_pack(rows, filename)

    if stat_pack.get("empty"):
        return {
            "stat_pack": stat_pack,
            "narrative": {
                "narrative": "The dataset has no rows, so there is nothing to narrate. Re-upload a file with data and try again.",
                "headline": "Empty dataset",
                "anchor_column": None,
            },
        }

    trimmed = dict(stat_pack)
    trimmed["columns"] = stat_pack["columns"][:40]

    user_prefs = {
        "audience": body.audience,
        "tone": body.tone,
        "length": body.length,
        "focus_column": body.focus_column,
    }
    user_message = (
        f"Author preferences:\n{json.dumps(user_prefs)}\n\n"
        f"Stat pack:\n{json.dumps(trimmed, default=str)}\n\n"
        "Write the narrative per the schema."
    )

    async with httpx.AsyncClient(timeout=90) as client:
        response = await client.post(
            ANTHROPIC_URL,
            headers={
                "x-api-key": ANTHROPIC_API_KEY,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": MODEL,
                "max_tokens": 1500,
                "system": NARRATIVE_SYSTEM_PROMPT,
                "messages": [{"role": "user", "content": user_message}],
            },
        )
        result = response.json()

    try:
        usage = result.get("usage", {}) or {}
        in_tok = int(usage.get("input_tokens", 0) or 0)
        out_tok = int(usage.get("output_tokens", 0) or 0)
        total = in_tok + out_tok
        if total > 0 and current_user.organisation_id:
            log_usage(
                db,
                organisation_id=current_user.organisation_id,
                user_id=current_user.id,
                kind="ai_tokens",
                quantity=total,
                cost_cents=_estimate_cost_cents(in_tok, out_tok),
                meta={"model": MODEL, "endpoint": "/ai/narrative", "input_tokens": in_tok, "output_tokens": out_tok, "file_id": file_id},
            )
    except Exception as exc:
        logger.warning("usage logging failed in /ai/narrative: %s", exc)

    content = result.get("content", [{}])
    text = content[0].get("text", "") if content else ""
    raw = text.strip()
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1] if "\n" in raw else raw
        if raw.rstrip().endswith("```"):
            raw = raw.rstrip()[:-3]
    try:
        narrative = json.loads(raw)
    except Exception:
        # Fallback: return the raw text so user isn't blocked.
        narrative = {
            "narrative": text,
            "headline": "AI narrative (unstructured)",
            "anchor_column": None,
            "raw": text,
        }

    return {"stat_pack": stat_pack, "narrative": narrative, "preferences": user_prefs}


# ---------------------------------------------------------------------------
# AI Auto-Report 2.0 — structured report + DOCX/PPTX export
# ---------------------------------------------------------------------------

REPORT_SYSTEM_PROMPT = """You are a senior data analyst writing a client-ready
analytical report. You will receive a statistical "stat pack" describing an
uploaded dataset. Produce a structured executive report.

Return ONLY raw JSON (no markdown). Schema:

{
  "title": "Title that names the dataset theme (5-10 words)",
  "executive_summary": "3-4 sentence exec summary grounded in specific numbers from the stat pack",
  "sections": [
    {
      "heading": "Section heading (e.g. 'Revenue composition')",
      "narrative": "2-4 sentences explaining what the data shows in business language, citing specific columns and numbers",
      "bullets": ["specific fact 1", "specific fact 2", "specific fact 3"]
    }
  ],
  "data_quality_notes": [
    "specific data quality observation 1",
    "specific data quality observation 2"
  ],
  "recommendations": [
    "concrete action 1",
    "concrete action 2",
    "concrete action 3"
  ]
}

Rules:
- 3-5 sections. Each bullet must reference a column name or a specific number.
- Do not invent columns or values that aren't in the stat pack.
- If data is low-n (<30 rows), add a caveat in the executive_summary."""


class ReportRequest(BaseModel):
    file_id: str


@router.post("/report/{file_id}")
async def ai_report(
    file_id: str,
    current_user: User = Depends(require_ai_enabled),
    db: Session = Depends(get_db),
):
    """Generate a structured analytical report from the file. Returns JSON that
    the frontend renders natively; the same JSON can be posted to /report/export
    to get a DOCX or PPTX."""
    if not ANTHROPIC_API_KEY:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured")
    _get_owned_file(file_id, current_user, db)
    rows, filename = load_file_data(file_id, current_user.organisation_id, db)
    stat_pack = build_stat_pack(rows, filename)

    if stat_pack.get("empty"):
        return {
            "stat_pack": stat_pack,
            "report": {
                "title": "Empty dataset",
                "executive_summary": "The file contains no data rows. Re-upload the source or verify the export.",
                "sections": [],
                "data_quality_notes": ["Dataset has 0 rows."],
                "recommendations": ["Re-export the data from the source system and upload again."],
            },
        }

    trimmed = dict(stat_pack)
    trimmed["columns"] = stat_pack["columns"][:40]
    user_message = (
        "Stat pack follows.\n\n"
        f"{json.dumps(trimmed, default=str)}\n\n"
        "Write the analytical report per the schema."
    )

    async with httpx.AsyncClient(timeout=120) as client:
        response = await client.post(
            ANTHROPIC_URL,
            headers={
                "x-api-key": ANTHROPIC_API_KEY,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": MODEL,
                "max_tokens": 3500,
                "system": REPORT_SYSTEM_PROMPT,
                "messages": [{"role": "user", "content": user_message}],
            },
        )
        result = response.json()

    try:
        usage = result.get("usage", {}) or {}
        in_tok = int(usage.get("input_tokens", 0) or 0)
        out_tok = int(usage.get("output_tokens", 0) or 0)
        total = in_tok + out_tok
        if total > 0 and current_user.organisation_id:
            log_usage(
                db,
                organisation_id=current_user.organisation_id,
                user_id=current_user.id,
                kind="ai_tokens",
                quantity=total,
                cost_cents=_estimate_cost_cents(in_tok, out_tok),
                meta={"model": MODEL, "endpoint": "/ai/report", "input_tokens": in_tok, "output_tokens": out_tok, "file_id": file_id},
            )
    except Exception as exc:
        logger.warning("usage logging failed in /ai/report: %s", exc)

    content = result.get("content", [{}])
    text = content[0].get("text", "") if content else ""
    raw = text.strip()
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1] if "\n" in raw else raw
        if raw.rstrip().endswith("```"):
            raw = raw.rstrip()[:-3]
    try:
        report = json.loads(raw)
    except Exception:
        report = {"title": "AI report", "executive_summary": text, "sections": [], "data_quality_notes": [], "recommendations": [], "raw": text}

    return {"stat_pack": stat_pack, "report": report}


class ReportExportRequest(BaseModel):
    file_id: str
    format: str  # "docx" | "pptx"
    report: dict  # the report JSON returned from /ai/report


BRAND_NAVY = (12, 20, 70)
BRAND_PINK = (233, 30, 140)
BRAND_TEAL = (0, 151, 178)


def _build_docx(report: dict, stat_pack: dict) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Title
    title_p = doc.add_paragraph()
    title_run = title_p.add_run(report.get("title", "Analytical Report"))
    title_run.bold = True
    title_run.font.size = Pt(22)
    title_run.font.color.rgb = RGBColor(*BRAND_NAVY)

    # Subtitle
    sub_p = doc.add_paragraph()
    sub_run = sub_p.add_run(
        f"{stat_pack.get('filename', '')} · {stat_pack.get('row_count', 0):,} rows · {stat_pack.get('column_count', 0)} columns"
    )
    sub_run.italic = True
    sub_run.font.size = Pt(10)
    sub_run.font.color.rgb = RGBColor(107, 114, 128)

    # Executive summary
    h1 = doc.add_paragraph()
    h1_run = h1.add_run("Executive summary")
    h1_run.bold = True
    h1_run.font.size = Pt(14)
    h1_run.font.color.rgb = RGBColor(*BRAND_PINK)
    doc.add_paragraph(report.get("executive_summary", ""))

    # Sections
    for sec in report.get("sections", []) or []:
        h = doc.add_paragraph()
        hr = h.add_run(sec.get("heading", "Section"))
        hr.bold = True
        hr.font.size = Pt(13)
        hr.font.color.rgb = RGBColor(*BRAND_NAVY)
        if sec.get("narrative"):
            doc.add_paragraph(sec["narrative"])
        for b in sec.get("bullets", []) or []:
            doc.add_paragraph(b, style="List Bullet")

    # Data quality
    if report.get("data_quality_notes"):
        h = doc.add_paragraph()
        hr = h.add_run("Data quality notes")
        hr.bold = True
        hr.font.size = Pt(13)
        hr.font.color.rgb = RGBColor(*BRAND_TEAL)
        for q in report["data_quality_notes"]:
            doc.add_paragraph(q, style="List Bullet")

    # Recommendations
    if report.get("recommendations"):
        h = doc.add_paragraph()
        hr = h.add_run("Recommended next steps")
        hr.bold = True
        hr.font.size = Pt(13)
        hr.font.color.rgb = RGBColor(*BRAND_PINK)
        for r in report["recommendations"]:
            doc.add_paragraph(r, style="List Number")

    # Stat pack appendix (columns table)
    doc.add_page_break()
    app = doc.add_paragraph()
    ar = app.add_run("Appendix: column statistics")
    ar.bold = True
    ar.font.size = Pt(13)
    ar.font.color.rgb = RGBColor(*BRAND_NAVY)

    cols = stat_pack.get("columns", []) or []
    if cols:
        table = doc.add_table(rows=1, cols=5)
        table.style = "Light Grid Accent 1"
        hdr_cells = table.rows[0].cells
        for i, h in enumerate(["Column", "Type", "Count", "Key stat", "Notes"]):
            hdr_cells[i].text = h
        for c in cols[:40]:
            row = table.add_row().cells
            row[0].text = str(c.get("column", ""))
            row[1].text = str(c.get("type", ""))
            row[2].text = str(c.get("count", ""))
            if c.get("type") == "numeric":
                row[3].text = f"mean={c.get('mean')}, min={c.get('min')}, max={c.get('max')}"
                notes = []
                if c.get("null_pct"): notes.append(f"{c['null_pct']}% null")
                if c.get("outliers_high") or c.get("outliers_low"):
                    notes.append(f"{(c.get('outliers_high') or 0) + (c.get('outliers_low') or 0)} outliers")
                row[4].text = ", ".join(notes)
            elif c.get("type") == "text":
                row[3].text = f"unique={c.get('unique')}, top={c.get('top_values', [{}])[0].get('value', '')}"
                row[4].text = f"top1 share {c.get('top1_share_pct', 0)}%"
            else:
                row[3].text = f"{c.get('min_date', '')} → {c.get('max_date', '')}"
                row[4].text = ""

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_pptx(report: dict, stat_pack: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_title_slide(text, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Navy background shape
        from pptx.shapes.autoshape import Shape
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(*BRAND_NAVY)
        bg.line.fill.background()
        # Title
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(2.3), Inches(12), Inches(2))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = text
        r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = RGBColor(255, 255, 255)
        # Subtitle
        sb = slide.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(12), Inches(1))
        sf = sb.text_frame
        sp = sf.paragraphs[0]
        sr = sp.add_run(); sr.text = subtitle
        sr.font.size = Pt(18); sr.font.color.rgb = RGBColor(233, 30, 140)

    def add_content_slide(heading, bullets, body=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Header bar
        bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.75))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*BRAND_NAVY)
        bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12), Inches(0.5))
        tf = tb.text_frame
        tp = tf.paragraphs[0]
        tr = tp.add_run(); tr.text = heading
        tr.font.size = Pt(24); tr.font.bold = True; tr.font.color.rgb = RGBColor(255, 255, 255)

        body_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.1), Inches(11.8), Inches(5.8))
        bf = body_box.text_frame; bf.word_wrap = True

        if body:
            p = bf.paragraphs[0]
            r = p.add_run(); r.text = body
            r.font.size = Pt(16); r.font.color.rgb = RGBColor(55, 65, 81)
            p.space_after = Pt(12)

        first = not bool(body)
        for b in bullets or []:
            p = bf.paragraphs[0] if first else bf.add_paragraph()
            first = False
            r = p.add_run(); r.text = "•  " + b
            r.font.size = Pt(15); r.font.color.rgb = RGBColor(55, 65, 81)
            p.space_after = Pt(8)

    # Title slide
    add_title_slide(
        report.get("title", "Analytical Report"),
        f"{stat_pack.get('filename', '')}  ·  {stat_pack.get('row_count', 0):,} rows  ·  {stat_pack.get('column_count', 0)} cols",
    )

    # Executive summary slide
    add_content_slide(
        "Executive summary",
        bullets=[],
        body=report.get("executive_summary", ""),
    )

    # Section slides
    for sec in report.get("sections", []) or []:
        add_content_slide(sec.get("heading", "Section"),
                          bullets=sec.get("bullets", []) or [],
                          body=sec.get("narrative", ""))

    # Data quality
    if report.get("data_quality_notes"):
        add_content_slide("Data quality", bullets=report["data_quality_notes"])

    # Recommendations
    if report.get("recommendations"):
        add_content_slide("Recommended next steps", bullets=report["recommendations"])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@router.post("/report/export")
async def ai_report_export(
    req: ReportExportRequest,
    current_user: User = Depends(require_ai_enabled),
    db: Session = Depends(get_db),
):
    """Turn a report JSON into a DOCX or PPTX file and stream it back."""
    file = _get_owned_file(req.file_id, current_user, db)
    # Re-derive a minimal stat_pack summary for the appendix and subtitle.
    # We don't want to recompute the full stat pack (expensive); use what
    # the file metadata gives us plus the columns from the report.
    stat_pack = {
        "filename": file.original_filename or file.filename,
        "row_count": file.row_count or 0,
        "column_count": file.column_count or 0,
        "columns": [],
    }
    if file.columns_json:
        try:
            cols = json.loads(file.columns_json) if isinstance(file.columns_json, str) else file.columns_json
            stat_pack["columns"] = [{"column": c, "type": "text", "count": file.row_count or 0} for c in cols]
        except Exception:
            pass

    fmt = (req.format or "docx").lower()
    safe_name = (file.original_filename or "report").rsplit(".", 1)[0]

    try:
        if fmt == "pptx":
            data = _build_pptx(req.report, stat_pack)
            return Response(
                content=data,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f'attachment; filename="{safe_name}_report.pptx"'},
            )
        else:
            data = _build_docx(req.report, stat_pack)
            return Response(
                content=data,
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={"Content-Disposition": f'attachment; filename="{safe_name}_report.docx"'},
            )
    except ModuleNotFoundError as e:
        raise HTTPException(
            status_code=500,
            detail=f"Report export requires {e.name}. Install python-docx and python-pptx.",
        )
